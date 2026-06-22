"""
RAG Engine with Evidence Tracking and Logging
Provides transparent visibility into what was retrieved and why
"""

import os
import warnings
from pathlib import Path
from typing import List, Dict, Optional, Tuple
from datetime import datetime
import ssl

# Fix for macOS SSL certificate verification errors when unstructured/spacy downloads models
try:
    _create_unverified_https_context = ssl._create_unverified_context
except AttributeError:
    pass
else:
    ssl._create_default_https_context = _create_unverified_https_context

import pdfplumber
from pptx import Presentation

from langchain_huggingface import HuggingFaceEmbeddings
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_core.documents import Document

from langchain_community.document_loaders import UnstructuredMarkdownLoader, WebBaseLoader
from langchain_community.retrievers import BM25Retriever
try:
    from langchain.retrievers import EnsembleRetriever
except ImportError:
    from langchain_classic.retrievers import EnsembleRetriever

try:
    from langchain_chroma import Chroma
except ImportError:
    from langchain_community.vectorstores import Chroma

from sentence_transformers import CrossEncoder

from data_models import Evidence


class RAGEngine:
    """RAG Engine with full transparency and evidence tracking"""

    def __init__(
        self,
        api_key: str,
        db_path: str = "./data/chroma_db",
        embed_model: str = "text-embedding-3-large",
        chunk_size: int = 2500,
        chunk_overlap: int = 400,
        verbose: bool = True
    ):
        self.api_key = api_key
        self.db_path = Path(db_path)
        self.db_path.mkdir(parents=True, exist_ok=True)

        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.verbose = verbose

        # Initialize Free Local Embeddings (Runs on your CPU/GPU)
        self.embeddings = HuggingFaceEmbeddings(
            model_name="all-MiniLM-L6-v2"
        )

        # Initialize vector store
        self.vectorstore = Chroma(
            collection_name="ask_my_docs_kb",
            embedding_function=self.embeddings,
            persist_directory=str(self.db_path)
        )
        
        # Initialize CrossEncoder for re-ranking
        self.reranker = CrossEncoder('cross-encoder/ms-marco-MiniLM-L-6-v2', max_length=512)

        self.indexing_log: List[Dict] = []
        self.retrieval_log: List[Dict] = []
        self.bm25_retriever = None
        self.all_documents = [] # Store for BM25

    def extract_pdf_text(self, pdf_path: Path) -> str:
        """Extract text from PDF with error handling"""
        if not pdf_path.exists():
            return ""

        try:
            with warnings.catch_warnings():
                warnings.simplefilter("ignore")

                with pdfplumber.open(pdf_path) as pdf:
                    text_parts = []
                    for page_num, page in enumerate(pdf.pages):
                        try:
                            text = page.extract_text()
                            if text and text.strip():
                                text_parts.append(f"[Page {page_num + 1}]\n{text.strip()}")
                        except Exception as e:
                            if self.verbose:
                                print(f"  ⚠️ Skipping page {page_num + 1}: {e}")
                            continue

                    return "\n\n".join(text_parts)
        except Exception as e:
            if self.verbose:
                print(f"  ❌ PDF extraction failed for {pdf_path.name}: {e}")
            return ""

    def extract_markdown_text(self, md_path: Path) -> str:
        """Extract text from Markdown files"""
        if not md_path.exists():
            return ""
        try:
            loader = UnstructuredMarkdownLoader(str(md_path))
            docs = loader.load()
            return "\n\n".join([doc.page_content for doc in docs])
        except Exception as e:
            if self.verbose:
                print(f"  ❌ Markdown extraction failed for {md_path.name}: {e}")
            return ""

    def extract_web_text(self, url: str) -> str:
        """Extract text from Web Pages"""
        try:
            loader = WebBaseLoader(url)
            docs = loader.load()
            return "\n\n".join([doc.page_content for doc in docs])
        except Exception as e:
            if self.verbose:
                print(f"  ❌ Web extraction failed for {url}: {e}")
            return ""

    def extract_pptx_text(self, pptx_path: Path) -> str:
        """Extract text from PowerPoint with structure preservation"""
        if not pptx_path.exists():
            return ""

        try:
            prs = Presentation(pptx_path)
            slide_texts = []

            for i, slide in enumerate(prs.slides):
                shapes_text = []
                for shape in slide.shapes:
                    try:
                        if hasattr(shape, "text") and shape.text.strip():
                            shapes_text.append(shape.text.strip())
                    except Exception:
                        continue

                if shapes_text:
                    slide_texts.append(f"[Slide {i + 1}]\n" + "\n".join(shapes_text))

            return "\n\n".join(slide_texts)
        except Exception as e:
            if self.verbose:
                print(f"  ❌ PPTX extraction failed for {pptx_path.name}: {e}")
            return ""

    def extract_approach_section(self, pptx_path: Path) -> str:
        """
        Extract ONLY the Approach section slides from a proposal

        This focuses on slides containing:
        - Timeline/schedule information
        - Governance structure (SteerCo, meetings)
        - Modules and activities
        - Workstream breakdowns

        These slides are weighted higher in RAG retrieval for learning patterns.
        """
        if not pptx_path.exists():
            return ""

        try:
            prs = Presentation(pptx_path)
            approach_slides = []

            approach_keywords = [
                'approach', 'methodology', 'timeline', 'workstream', 'work stream',
                'governance', 'steerco', 'steering', 'module', 'activity',
                'deliverable', 'week', 'phase', 'gantt', 'schedule'
            ]

            for i, slide in enumerate(prs.slides):
                shapes_text = []
                all_slide_text = ""

                for shape in slide.shapes:
                    try:
                        if hasattr(shape, "text") and shape.text.strip():
                            text = shape.text.strip()
                            shapes_text.append(text)
                            all_slide_text += text.lower() + " "
                    except Exception:
                        continue

                # Check if this slide is part of approach section
                is_approach_slide = any(keyword in all_slide_text for keyword in approach_keywords)

                if is_approach_slide and shapes_text:
                    slide_content = f"[Approach Slide {i + 1}]\n" + "\n".join(shapes_text)
                    approach_slides.append(slide_content)

            if self.verbose and approach_slides:
                print(f"  🎯 Found {len(approach_slides)} approach slides in {pptx_path.name}")

            return "\n\n".join(approach_slides)

        except Exception as e:
            if self.verbose:
                print(f"  ❌ Approach extraction failed for {pptx_path.name}: {e}")
            return ""

    def chunk_document(
        self,
        text: str,
        doc_type: str,
        source: str,
        preserve_structure: bool = True
    ) -> List[Document]:
        """Create intelligent document chunks"""
        if not text.strip():
            return []

        chunks = []

        if doc_type == "proposal" and preserve_structure:
            # Preserve slide structure for proposals
            slides = text.split("[Slide ")

            for slide in slides:
                if not slide.strip():
                    continue

                # Extract slide number and content
                slide_parts = slide.split("]", 1)
                slide_num = slide_parts[0] if len(slide_parts) > 1 else "0"
                slide_content = slide_parts[1] if len(slide_parts) > 1 else slide

                if len(slide_content) > self.chunk_size:
                    # Split large slides by paragraphs
                    paragraphs = slide_content.split('\n\n')
                    current_chunk = ""

                    for para in paragraphs:
                        if len(current_chunk + para) > self.chunk_size and current_chunk:
                            chunks.append(Document(
                                page_content=current_chunk.strip(),
                                metadata={
                                    "source": source,
                                    "doc_type": doc_type,
                                    "slide": slide_num,
                                    "chunk_type": "slide_section"
                                }
                            ))
                            current_chunk = para
                        else:
                            current_chunk += "\n\n" + para if current_chunk else para

                    if current_chunk:
                        chunks.append(Document(
                            page_content=current_chunk.strip(),
                            metadata={
                                "source": source,
                                "doc_type": doc_type,
                                "slide": slide_num,
                                "chunk_type": "slide_section"
                            }
                        ))
                else:
                    chunks.append(Document(
                        page_content=slide_content.strip(),
                        metadata={
                            "source": source,
                            "doc_type": doc_type,
                            "slide": slide_num,
                            "chunk_type": "full_slide"
                        }
                    ))

        else:
            # Document chunking with semantic awareness
            splitter = RecursiveCharacterTextSplitter(
                chunk_size=self.chunk_size,
                chunk_overlap=self.chunk_overlap,
                separators=["\n\n", "\n", ". ", "! ", "? ", " "],
                length_function=len,
            )

            texts = splitter.split_text(text)

            for i, chunk in enumerate(texts):
                if chunk.strip():
                    chunks.append(Document(
                        page_content=chunk.strip(),
                        metadata={
                            "source": source,
                            "doc_type": doc_type,
                            "chunk": i,
                            "chunk_type": "semantic"
                        }
                    ))

        return chunks

    def index_documents(
        self,
        rfp_dir: str = "./data/rfps",
        proposal_dir: str = "./data/proposals"
    ) -> Dict[str, any]:
        """Index all documents with detailed logging"""
        rfp_path = Path(rfp_dir)
        proposal_path = Path(proposal_dir)

        indexed_count = 0
        errors = []
        file_details = []

        start_time = datetime.now()

        # Index RFPs
        if self.verbose:
            print("📄 Indexing RFP documents...")

        rfp_files = list(rfp_path.glob("*.pdf"))
        for rfp_file in rfp_files:
            try:
                if self.verbose:
                    print(f"  Processing: {rfp_file.name}")

                text = self.extract_pdf_text(rfp_file)
                if not text:
                    errors.append(f"No text extracted from {rfp_file.name}")
                    continue

                docs = self.chunk_document(text, "rfp", rfp_file.name)
                if docs:
                    self.vectorstore.add_documents(docs)
                    indexed_count += len(docs)

                    file_details.append({
                        "file": rfp_file.name,
                        "type": "rfp",
                        "chunks": len(docs),
                        "characters": len(text)
                    })

                    if self.verbose:
                        print(f"    ✅ Indexed {len(docs)} chunks ({len(text):,} chars)")
                else:
                    errors.append(f"No chunks created from {rfp_file.name}")

            except Exception as e:
                error_msg = f"Failed to index {rfp_file.name}: {str(e)}"
                errors.append(error_msg)
                if self.verbose:
                    print(f"    ❌ {error_msg}")

        # Index Proposals
        if self.verbose:
            print("\n🎨 Indexing Proposal documents...")

        proposal_files = list(proposal_path.glob("*.pptx"))
        for proposal_file in proposal_files:
            try:
                if self.verbose:
                    print(f"  Processing: {proposal_file.name}")

                # Extract ALL slides
                text = self.extract_pptx_text(proposal_file)
                if not text:
                    errors.append(f"No text extracted from {proposal_file.name}")
                    continue

                docs = self.chunk_document(text, "proposal", proposal_file.name)

                # ALSO extract APPROACH section specifically with higher weight
                approach_text = self.extract_approach_section(proposal_file)
                approach_docs = []

                if approach_text:
                    approach_docs = self.chunk_document(approach_text, "proposal_approach", proposal_file.name)

                    # Add metadata flag to approach chunks for prioritized retrieval
                    for doc in approach_docs:
                        doc.metadata['is_approach'] = True
                        doc.metadata['priority'] = 'high'

                # Add all documents to vectorstore
                all_docs = docs + approach_docs

                if all_docs:
                    self.vectorstore.add_documents(all_docs)
                    indexed_count += len(all_docs)

                    file_details.append({
                        "file": proposal_file.name,
                        "type": "proposal",
                        "chunks": len(docs),
                        "approach_chunks": len(approach_docs),
                        "total_chunks": len(all_docs),
                        "characters": len(text),
                        "approach_chars": len(approach_text) if approach_text else 0
                    })

                    if self.verbose:
                        print(f"    ✅ Indexed {len(docs)} general chunks + {len(approach_docs)} approach chunks")
                        print(f"       Total: {len(all_docs)} chunks ({len(text):,} chars)")
                else:
                    errors.append(f"No chunks created from {proposal_file.name}")

            except Exception as e:
                error_msg = f"Failed to index {proposal_file.name}: {str(e)}"
                errors.append(error_msg)
                if self.verbose:
                    print(f"    ❌ {error_msg}")
                    
        # Index Markdown files
        md_dir = Path("./data/markdown")
        if md_dir.exists():
            if self.verbose:
                print("\n📝 Indexing Markdown documents...")
                
            md_files = list(md_dir.glob("*.md"))
            for md_file in md_files:
                try:
                    if self.verbose:
                        print(f"  Processing: {md_file.name}")

                    text = self.extract_markdown_text(md_file)
                    if not text:
                        errors.append(f"No text extracted from {md_file.name}")
                        continue

                    docs = self.chunk_document(text, "markdown", md_file.name, preserve_structure=False)
                    if docs:
                        self.vectorstore.add_documents(docs)
                        self.all_documents.extend(docs)
                        indexed_count += len(docs)

                        file_details.append({
                            "file": md_file.name,
                            "type": "markdown",
                            "chunks": len(docs),
                            "characters": len(text)
                        })

                        if self.verbose:
                            print(f"    ✅ Indexed {len(docs)} chunks ({len(text):,} chars)")
                    else:
                        errors.append(f"No chunks created from {md_file.name}")

                except Exception as e:
                    error_msg = f"Failed to index {md_file.name}: {str(e)}"
                    errors.append(error_msg)
                    if self.verbose:
                        print(f"    ❌ {error_msg}")
                        
        # Setup BM25 Retriever if documents exist
        if self.all_documents:
            self.bm25_retriever = BM25Retriever.from_documents(self.all_documents)
            self.bm25_retriever.k = 10

        duration = (datetime.now() - start_time).total_seconds()

        # Log indexing results
        log_entry = {
            "timestamp": datetime.now(),
            "total_chunks": indexed_count,
            "rfp_count": len([f for f in file_details if f["type"] == "rfp"]),
            "proposal_count": len([f for f in file_details if f["type"] == "proposal"]),
            "errors": errors,
            "duration_seconds": duration,
            "file_details": file_details
        }

        self.indexing_log.append(log_entry)

        if self.verbose:
            print(f"\n{'='*60}")
            print(f"✅ Indexing Complete!")
            print(f"   Total chunks indexed: {indexed_count}")
            print(f"   RFP files: {len([f for f in file_details if f['type'] == 'rfp'])}")
            print(f"   Proposal files: {len([f for f in file_details if f['type'] == 'proposal'])}")
            print(f"   Duration: {duration:.2f}s")
            if errors:
                print(f"   ⚠️  Errors: {len(errors)}")
            print(f"{'='*60}\n")

        return log_entry

    def retrieve_evidence(
        self,
        query: str,
        doc_types: Optional[List[str]] = None,
        k: int = 15,
        enhance_query: bool = True
    ) -> List[Evidence]:
        """Hybrid Retrieval (BM25 + Vector) and Cross-Encoder Re-ranking"""
        start_time = datetime.now()

        # Enhance query with domain-specific terms (Optional based on domain)
        if enhance_query:
            enhanced_query = f"{query}"
        else:
            enhanced_query = query

        if self.verbose:
            print(f"\n🔍 Hybrid RAG Retrieval & Re-ranking")
            print(f"   Query: {query}")
            print(f"   Doc types: {doc_types or 'all'}")
            print(f"   Top-K final: {k}")

        try:
            # Use EnsembleRetriever for Hybrid Search if BM25 is configured
            if self.bm25_retriever:
                vector_retriever = self.vectorstore.as_retriever(search_kwargs={"k": k*2})
                ensemble_retriever = EnsembleRetriever(
                    retrievers=[self.bm25_retriever, vector_retriever],
                    weights=[0.3, 0.7]
                )
                retrieved_docs = ensemble_retriever.invoke(enhanced_query)
            else:
                retrieved_docs = self.vectorstore.similarity_search(enhanced_query, k=k*2)
                
            if not retrieved_docs:
                return []
                
            # Filter by doc_types if specified
            if doc_types:
                retrieved_docs = [doc for doc in retrieved_docs if doc.metadata.get("doc_type") in doc_types]

            # Cross-Encoder Re-ranking
            pairs = [[query, doc.page_content] for doc in retrieved_docs]
            scores = self.reranker.predict(pairs)
            
            # Combine docs and scores, then sort
            doc_score_pairs = list(zip(retrieved_docs, scores))
            doc_score_pairs.sort(key=lambda x: x[1], reverse=True) # Higher score is better for CrossEncoder
            
            # Take top K
            top_scored_docs = doc_score_pairs[:k]

            # Convert to Evidence objects
            evidence_list = []
            for doc, score in top_scored_docs:
                evidence = Evidence(
                    source_file=doc.metadata.get("source", "unknown"),
                    doc_type=doc.metadata.get("doc_type", "unknown"),
                    content=doc.page_content,
                    similarity_score=float(score), # Now represents cross-encoder score
                    metadata=doc.metadata
                )
                evidence_list.append(evidence)

            duration = (datetime.now() - start_time).total_seconds()

            # Log retrieval
            log_entry = {
                "timestamp": datetime.now(),
                "query": query,
                "enhanced_query": enhanced_query if enhance_query else None,
                "doc_types": doc_types,
                "k": k,
                "results_count": len(evidence_list),
                "duration_seconds": duration,
                "top_sources": list(set([e.source_file for e in evidence_list[:5]]))
            }

            self.retrieval_log.append(log_entry)

            if self.verbose:
                print(f"\n   ✅ Retrieved {len(evidence_list)} results in {duration:.2f}s")
                print(f"   Top sources:")
                for source in log_entry["top_sources"]:
                    source_count = len([e for e in evidence_list if e.source_file == source])
                    avg_score = sum([e.similarity_score for e in evidence_list if e.source_file == source]) / source_count
                    print(f"      • {source}: {source_count} chunks (avg similarity: {avg_score:.3f})")

            return evidence_list

        except Exception as e:
            if self.verbose:
                print(f"   ❌ Retrieval failed: {e}")
            return []

    def retrieve_approach_evidence(
        self,
        query: str,
        k: int = 20,
        approach_only: bool = True
    ) -> List[Evidence]:
        """
        Retrieve evidence specifically from APPROACH sections of proposals

        This method prioritizes approach slides over general proposal content,
        ensuring the AI learns from actual approach section patterns.

        Args:
            query: Search query
            k: Number of results to retrieve
            approach_only: If True, only return approach section chunks

        Returns:
            List of Evidence objects from approach sections
        """
        start_time = datetime.now()

        # Enhance query for approach-specific retrieval
        enhanced_query = f"{query} approach methodology timeline workstream module activity governance"

        if self.verbose:
            print(f"\n🎯 APPROACH-FOCUSED RAG Retrieval")
            print(f"   Query: {query}")
            print(f"   Enhanced: {enhanced_query}")
            print(f"   Approach only: {approach_only}")
            print(f"   Top-K: {k}")

        try:
            # Retrieve more results than needed for filtering
            retrieve_k = k * 3 if approach_only else k

            docs_scores = self.vectorstore.similarity_search_with_score(enhanced_query, k=retrieve_k)

            # Filter and prioritize approach sections
            approach_evidence = []
            general_evidence = []

            for doc, score in docs_scores:
                evidence = Evidence(
                    source_file=doc.metadata.get("source", "unknown"),
                    doc_type=doc.metadata.get("doc_type", "unknown"),
                    content=doc.page_content,
                    similarity_score=float(score),
                    metadata=doc.metadata
                )

                # Check if this is an approach section chunk
                if doc.metadata.get('is_approach', False):
                    # Boost approach section results (lower score = better)
                    evidence.similarity_score *= 0.7  # 30% boost
                    approach_evidence.append(evidence)
                else:
                    general_evidence.append(evidence)

            # Combine: prioritize approach sections
            if approach_only:
                evidence_list = approach_evidence[:k]
            else:
                # Mix approach (70%) with general (30%)
                approach_count = int(k * 0.7)
                general_count = k - approach_count
                evidence_list = approach_evidence[:approach_count] + general_evidence[:general_count]

            # Sort by relevance
            evidence_list.sort(key=lambda x: x.similarity_score)

            duration = (datetime.now() - start_time).total_seconds()

            # Log retrieval
            approach_chunk_count = len([e for e in evidence_list if e.metadata.get('is_approach')])

            if self.verbose:
                print(f"\n   ✅ Retrieved {len(evidence_list)} results ({approach_chunk_count} from approach sections)")
                print(f"   Duration: {duration:.2f}s")
                print(f"   Top approach sources:")

                approach_sources = {}
                for e in evidence_list:
                    if e.metadata.get('is_approach'):
                        approach_sources[e.source_file] = approach_sources.get(e.source_file, 0) + 1

                for source, count in sorted(approach_sources.items(), key=lambda x: x[1], reverse=True)[:5]:
                    print(f"      • {source}: {count} approach chunks")

            return evidence_list

        except Exception as e:
            if self.verbose:
                print(f"   ❌ Approach retrieval failed: {e}")
            return []

    def get_collection_stats(self) -> Dict:
        """Get statistics about the vector store collection"""
        try:
            collection = self.vectorstore._collection
            count = collection.count()

            # Get all documents to analyze
            all_docs = self.vectorstore.get()

            sources = {}
            doc_types = {"rfp": 0, "proposal": 0}

            if all_docs and 'metadatas' in all_docs:
                for metadata in all_docs['metadatas']:
                    source = metadata.get('source', 'unknown')
                    doc_type = metadata.get('doc_type', 'unknown')

                    sources[source] = sources.get(source, 0) + 1
                    if doc_type in doc_types:
                        doc_types[doc_type] += 1

            return {
                "total_chunks": count,
                "unique_sources": len(sources),
                "doc_types": doc_types,
                "sources": sources
            }
        except Exception as e:
            if self.verbose:
                print(f"⚠️ Could not get collection stats: {e}")
            return {"total_chunks": 0, "unique_sources": 0}

    def clear_collection(self):
        """Clear all documents from the vector store"""
        try:
            # Delete and recreate collection
            self.vectorstore._client.delete_collection(self.vectorstore._collection.name)

            self.vectorstore = Chroma(
                collection_name="rfp_proposal_kb",
                embedding_function=self.embeddings,
                persist_directory=str(self.db_path)
            )

            if self.verbose:
                print("✅ Collection cleared")
        except Exception as e:
            if self.verbose:
                print(f"❌ Failed to clear collection: {e}")
